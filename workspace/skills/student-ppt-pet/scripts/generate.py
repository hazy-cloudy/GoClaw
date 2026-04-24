import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
DEFAULT_FONT = "Microsoft YaHei"
TITLE_FONT = "Microsoft YaHei UI"

THEME = {
    "navy": "0F2747",
    "blue": "1F5EFF",
    "sky": "EAF2FF",
    "cyan": "4DA8FF",
    "ink": "1E2A39",
    "gray": "5E6B78",
    "light": "F7F9FC",
    "white": "FFFFFF",
    "border": "D8E1F0",
    "accent": "FF8A3D",
    "green": "1FA971",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.replace("#", "")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def normalize_lines(items):
    if not isinstance(items, list):
        return []
    return [str(item).strip() for item in items if str(item).strip()]


def ensure_slide_size(prs: Presentation) -> None:
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H


def set_background(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_box(slide, left, top, width, height, fill_color, line_color=None, radius=True):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    shape.line.color.rgb = rgb(line_color or fill_color)
    return shape


def add_line(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    return shape


def style_text_frame(text_frame, *, font_name=DEFAULT_FONT, font_size=20, color="1E2A39",
                     bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    text_frame.vertical_anchor = valign
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)


def add_text(slide, left, top, width, height, text, *,
             font_name=DEFAULT_FONT, font_size=18, color="1E2A39",
             bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    style_text_frame(
        tf,
        font_name=font_name,
        font_size=font_size,
        color=color,
        bold=bold,
        align=align,
        valign=valign,
    )
    return textbox


def add_bullets(slide, left, top, width, height, bullets, *,
                font_size=18, color="1E2A39", level=0, max_items=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    items = bullets[:max_items] if bullets else ["待补充"]
    first = tf.paragraphs[0]
    first.text = items[0]
    first.level = level
    first.space_after = Pt(10)
    style_text_frame(tf, font_size=font_size, color=color)

    for item in items[1:]:
        p = tf.add_paragraph()
        p.text = item
        p.level = level
        p.space_after = Pt(10)
        for run in p.runs:
            run.font.name = DEFAULT_FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = rgb(color)
    return box


def add_footer(slide, page_num: int):
    add_line(slide, Inches(0.6), Inches(7.05), Inches(12.1), Inches(0.03), THEME["border"])
    add_text(
        slide,
        Inches(0.7),
        Inches(7.08),
        Inches(4.5),
        Inches(0.25),
        "Student PPT Pet · Academic Presentation",
        font_size=9,
        color=THEME["gray"],
    )
    badge = add_box(slide, Inches(12.05), Inches(6.93), Inches(0.55), Inches(0.34), THEME["navy"])
    add_text(
        slide,
        badge.left,
        badge.top + Inches(0.01),
        badge.width,
        badge.height,
        str(page_num),
        font_name=TITLE_FONT,
        font_size=12,
        color=THEME["white"],
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def infer_layout(slide_data: dict) -> str:
    layout = str(slide_data.get("layout", "")).lower().strip()
    if layout:
        return layout

    slide_type = str(slide_data.get("type", "content")).lower()
    if slide_type == "title":
        return "title"

    title = str(slide_data.get("title", ""))
    if "目录" in title:
        return "agenda"
    if any(k in title for k in ["方法", "设计", "流程", "架构", "实验设置"]):
        return "process"
    if any(k in title for k in ["结果", "分析", "数据", "挑战", "对比"]):
        return "results"
    if any(k in title for k in ["总结", "展望", "结论", "感谢"]):
        return "summary"
    return "content"


def add_cover_slide(prs: Presentation, slide_data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME["navy"])

    add_box(slide, Inches(0.8), Inches(0.75), Inches(0.18), Inches(5.8), THEME["accent"], radius=False)
    add_box(slide, Inches(9.7), Inches(0.8), Inches(2.2), Inches(2.2), THEME["blue"])
    add_box(slide, Inches(10.35), Inches(1.45), Inches(1.2), Inches(1.2), THEME["cyan"])
    add_box(slide, Inches(8.85), Inches(4.95), Inches(2.7), Inches(1.05), "214C99")

    add_text(
        slide,
        Inches(1.35),
        Inches(1.2),
        Inches(7.2),
        Inches(1.4),
        slide_data.get("title", "Academic Presentation"),
        font_name=TITLE_FONT,
        font_size=26,
        color=THEME["white"],
        bold=True,
    )
    subtitle = "\n".join(normalize_lines(slide_data.get("content", [])))
    add_text(
        slide,
        Inches(1.38),
        Inches(2.75),
        Inches(6.4),
        Inches(1.5),
        subtitle or "课程汇报 / 学术答辩 / 研究展示",
        font_size=15,
        color="D6E2FF",
    )
    add_text(
        slide,
        Inches(1.38),
        Inches(5.75),
        Inches(4.2),
        Inches(0.4),
        "Student PPT Pet",
        font_size=11,
        color="B8CAE8",
    )
    add_text(
        slide,
        Inches(9.05),
        Inches(5.18),
        Inches(2.2),
        Inches(0.3),
        "Structured • Academic • Editable",
        font_size=10,
        color=THEME["white"],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return slide


def add_agenda_slide(prs: Presentation, slide_data: dict, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME["light"])
    add_box(slide, Inches(0.45), Inches(0.45), Inches(12.25), Inches(0.55), THEME["navy"])
    add_text(slide, Inches(0.75), Inches(0.56), Inches(4), Inches(0.25), slide_data.get("title", "目录"),
             font_name=TITLE_FONT, font_size=24, color=THEME["white"], bold=True)

    items = normalize_lines(slide_data.get("content", []))
    positions = [
        (Inches(0.9), Inches(1.5)),
        (Inches(6.7), Inches(1.5)),
        (Inches(0.9), Inches(3.15)),
        (Inches(6.7), Inches(3.15)),
        (Inches(0.9), Inches(4.8)),
        (Inches(6.7), Inches(4.8)),
    ]
    for idx, item in enumerate(items[:6], start=1):
        left, top = positions[idx - 1]
        add_box(slide, left, top, Inches(5.0), Inches(1.2), THEME["white"], THEME["border"])
        circle = add_box(slide, left + Inches(0.18), top + Inches(0.16), Inches(0.75), Inches(0.75), THEME["blue"])
        add_text(slide, circle.left, circle.top + Inches(0.01), circle.width, circle.height, f"{idx:02d}",
                 font_name=TITLE_FONT, font_size=15, color=THEME["white"], bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, left + Inches(1.1), top + Inches(0.18), Inches(3.55), Inches(0.7), item,
                 font_name=DEFAULT_FONT, font_size=18, color=THEME["ink"], bold=True, valign=MSO_ANCHOR.MIDDLE)
    add_footer(slide, page_num)
    return slide


def add_process_slide(prs: Presentation, slide_data: dict, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME["white"])
    add_text(slide, Inches(0.7), Inches(0.55), Inches(8.5), Inches(0.5), slide_data.get("title", ""),
             font_name=TITLE_FONT, font_size=24, color=THEME["navy"], bold=True)
    add_line(slide, Inches(0.72), Inches(1.1), Inches(2.6), Inches(0.05), THEME["accent"])

    bullets = normalize_lines(slide_data.get("content", []))
    steps = bullets[:4] if bullets else ["步骤一", "步骤二", "步骤三", "步骤四"]
    step_width = Inches(2.75)
    step_height = Inches(2.35)
    base_top = Inches(2.0)
    for idx, step in enumerate(steps):
        left = Inches(0.75) + idx * Inches(3.05)
        card = add_box(slide, left, base_top, step_width, step_height, THEME["sky"], THEME["border"])
        badge = add_box(slide, left + Inches(0.22), base_top + Inches(0.2), Inches(0.52), Inches(0.52), THEME["blue"])
        add_text(slide, badge.left, badge.top, badge.width, badge.height, str(idx + 1),
                 font_name=TITLE_FONT, font_size=14, color=THEME["white"], bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, left + Inches(0.22), base_top + Inches(0.9), Inches(2.2), Inches(1.0), step,
                 font_size=17, color=THEME["ink"], bold=True)
        if idx < len(steps) - 1:
            add_line(slide, left + Inches(2.8), base_top + Inches(1.12), Inches(0.32), Inches(0.04), THEME["cyan"])
    add_footer(slide, page_num)
    return slide


def add_results_chart(slide, chart_data_dict, left, top, width, height):
    chart_data = ChartData()
    categories = chart_data_dict.get("categories", [])
    values = chart_data_dict.get("values", [])
    chart_data.categories = categories
    chart_data.add_series(chart_data_dict.get("series_name", "结果"), values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left,
        top,
        width,
        height,
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.tick_labels.font.size = Pt(11)
    chart.value_axis.tick_labels.font.size = Pt(11)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(10)
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = rgb(THEME["blue"])


def add_results_slide(prs: Presentation, slide_data: dict, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, "F5F8FE")
    add_text(slide, Inches(0.7), Inches(0.55), Inches(8.5), Inches(0.45), slide_data.get("title", ""),
             font_name=TITLE_FONT, font_size=24, color=THEME["navy"], bold=True)
    add_line(slide, Inches(0.72), Inches(1.08), Inches(2.9), Inches(0.05), THEME["green"])

    add_box(slide, Inches(0.72), Inches(1.55), Inches(5.2), Inches(4.95), THEME["white"], THEME["border"])
    add_text(slide, Inches(1.0), Inches(1.85), Inches(4.5), Inches(0.3), "核心发现", font_size=16, color=THEME["blue"], bold=True)
    add_bullets(slide, Inches(1.0), Inches(2.2), Inches(4.45), Inches(3.7), normalize_lines(slide_data.get("content", [])),
                font_size=17, color=THEME["ink"])

    metrics = slide_data.get("metrics", [])
    if isinstance(metrics, list) and metrics:
        for idx, metric in enumerate(metrics[:3]):
            left = Inches(6.2) + idx * Inches(2.05)
            add_box(slide, left, Inches(1.75), Inches(1.75), Inches(1.25), THEME["navy"])
            add_text(slide, left + Inches(0.12), Inches(1.95), Inches(1.5), Inches(0.35),
                     str(metric.get("value", "--")), font_name=TITLE_FONT, font_size=20, color=THEME["white"], bold=True,
                     align=PP_ALIGN.CENTER)
            add_text(slide, left + Inches(0.1), Inches(2.55), Inches(1.55), Inches(0.35),
                     str(metric.get("label", "")), font_size=11, color="DDE7FA",
                     align=PP_ALIGN.CENTER)
    else:
        cards = normalize_lines(slide_data.get("content", []))[:3]
        for idx, item in enumerate(cards):
            left = Inches(6.2)
            top = Inches(1.75) + idx * Inches(1.45)
            add_box(slide, left, top, Inches(5.4), Inches(1.15), THEME["white"], THEME["border"])
            add_line(slide, left, top, Inches(0.12), Inches(1.15), [THEME["blue"], THEME["green"], THEME["accent"]][idx % 3])
            add_text(slide, left + Inches(0.28), top + Inches(0.18), Inches(4.8), Inches(0.65), item,
                     font_size=16, color=THEME["ink"], bold=True, valign=MSO_ANCHOR.MIDDLE)

    chart = slide_data.get("chart", {})
    if isinstance(chart, dict) and chart.get("categories") and chart.get("values"):
        add_results_chart(slide, chart, Inches(6.3), Inches(4.55), Inches(5.2), Inches(1.5))

    add_footer(slide, page_num)
    return slide


def add_summary_slide(prs: Presentation, slide_data: dict, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME["navy"])
    add_text(slide, Inches(0.82), Inches(0.62), Inches(5.5), Inches(0.5), slide_data.get("title", ""),
             font_name=TITLE_FONT, font_size=24, color=THEME["white"], bold=True)
    add_text(slide, Inches(0.82), Inches(1.32), Inches(4.6), Inches(0.35), "Key Takeaways", font_size=13, color="B8CAE8")

    add_box(slide, Inches(0.82), Inches(1.85), Inches(7.1), Inches(4.65), "17365F")
    bullets = normalize_lines(slide_data.get("content", []))
    for idx, item in enumerate(bullets[:5]):
        top = Inches(2.15) + idx * Inches(0.82)
        badge = add_box(slide, Inches(1.12), top, Inches(0.32), Inches(0.32), THEME["accent"])
        add_text(slide, badge.left, badge.top - Inches(0.01), badge.width, badge.height, "✓",
                 font_size=11, color=THEME["white"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(1.58), top - Inches(0.02), Inches(5.7), Inches(0.42), item,
                 font_size=17, color=THEME["white"], bold=False, valign=MSO_ANCHOR.MIDDLE)

    add_box(slide, Inches(8.35), Inches(1.85), Inches(3.95), Inches(4.65), "214C99")
    add_text(slide, Inches(8.7), Inches(2.2), Inches(3.2), Inches(0.4), "答辩提示", font_name=TITLE_FONT,
             font_size=18, color=THEME["white"], bold=True)
    hints = [
        "先讲问题，再讲方法，最后讲结果。",
        "指标解释要简洁，避免堆砌术语。",
        "结论页要明确贡献与下一步。"
    ]
    add_bullets(slide, Inches(8.68), Inches(2.8), Inches(2.95), Inches(2.8), hints, font_size=15, color=THEME["white"])
    add_footer(slide, page_num)
    return slide


def add_content_slide(prs: Presentation, slide_data: dict, page_num: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, THEME["light"])
    add_box(slide, Inches(0.65), Inches(0.7), Inches(2.55), Inches(5.95), THEME["navy"])
    add_text(slide, Inches(0.92), Inches(1.2), Inches(2.0), Inches(1.5), slide_data.get("title", ""),
             font_name=TITLE_FONT, font_size=22, color=THEME["white"], bold=True)
    add_text(slide, Inches(0.95), Inches(5.7), Inches(1.8), Inches(0.35), "Academic Section", font_size=10, color="B8CAE8")

    add_box(slide, Inches(3.45), Inches(0.95), Inches(8.95), Inches(5.45), THEME["white"], THEME["border"])
    add_text(slide, Inches(3.82), Inches(1.28), Inches(3.5), Inches(0.28), "核心内容", font_size=15, color=THEME["blue"], bold=True)
    add_bullets(slide, Inches(3.82), Inches(1.7), Inches(5.65), Inches(4.2), normalize_lines(slide_data.get("content", [])),
                font_size=18, color=THEME["ink"])

    bullets = normalize_lines(slide_data.get("content", []))
    right_cards = bullets[:2] if bullets else ["结论要点一", "结论要点二"]
    for idx, item in enumerate(right_cards):
        top = Inches(1.72) + idx * Inches(2.05)
        add_box(slide, Inches(9.75), top, Inches(2.25), Inches(1.55), THEME["sky"], THEME["border"])
        add_text(slide, Inches(9.96), top + Inches(0.18), Inches(1.8), Inches(1.1), item,
                 font_size=14, color=THEME["ink"], bold=True, valign=MSO_ANCHOR.MIDDLE)

    add_footer(slide, page_num)
    return slide


def build_slide(prs: Presentation, slide_data: dict, page_num: int):
    layout = infer_layout(slide_data)
    if layout == "title":
        return add_cover_slide(prs, slide_data)
    if layout == "agenda":
        return add_agenda_slide(prs, slide_data, page_num)
    if layout == "process":
        return add_process_slide(prs, slide_data, page_num)
    if layout == "results":
        return add_results_slide(prs, slide_data, page_num)
    if layout in {"summary", "closing"}:
        return add_summary_slide(prs, slide_data, page_num)
    return add_content_slide(prs, slide_data, page_num)


def build_presentation(plan: dict) -> Presentation:
    prs = Presentation()
    ensure_slide_size(prs)
    prs.core_properties.title = str(plan.get("title", "Academic Presentation"))
    prs.core_properties.author = str(plan.get("author", "Unknown Author"))

    slides = plan.get("slides", [])
    if not isinstance(slides, list) or not slides:
        raise ValueError("slides must be a non-empty list")

    for idx, slide_data in enumerate(slides, start=1):
        if not isinstance(slide_data, dict):
            continue
        build_slide(prs, slide_data, idx)
    return prs


def create_presentation(plan_file: str, output_file: str) -> str:
    plan_path = Path(plan_file).expanduser().resolve()
    output_path = Path(output_file).expanduser().resolve()

    if not plan_path.exists():
        raise FileNotFoundError(f"plan file not found: {plan_path}")

    output_path.parent.mkdir(parents=True, exist_ok=True)

    with plan_path.open("r", encoding="utf-8") as f:
        plan = json.load(f)

    prs = build_presentation(plan)
    prs.save(str(output_path))
    return str(output_path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate a styled PPTX file from a JSON plan")
    parser.add_argument("--plan-file", required=True, help="Path to JSON plan")
    parser.add_argument("--output-file", required=True, help="Destination PPTX path")
    args = parser.parse_args()

    try:
        output_path = create_presentation(args.plan_file, args.output_file)
    except Exception as exc:
        print(f"[Error] {exc}")
        return 1

    print(f"[Success] Presentation successfully generated at: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
